from __future__ import annotations

from pptx import Presentation

from app.domain.enums import SlideSize
from app.domain.sections import CoverSection
from app.domain.style import (
    BlockLayout,
    EdgeInsets,
    SectionStyle,
    TextBlockStyle,
    TextStyle,
)
from app.pptx.builder import build_pptx
from app.pptx.layout import Rect, resolve_block_rect
from app.services.generation import SlideModel


def _style(block_id: str, anchor: str, **margins: float) -> dict:
    return {
        "margin": 0.8,
        "blocks": {
            block_id: {
                "layout": {
                    "anchor": anchor,
                    "margin": margins,
                }
            }
        },
    }


def test_unconfigured_block_preserves_legacy_rectangle():
    base = Rect(0.8, 2.4, 11.733, 2.0)
    assert resolve_block_rect(
        base,
        block_id="title",
        default_anchor="middle_center",
        style={"margin": 0.8},
        slide_width=13.333,
        slide_height=7.5,
    ) == base


def test_all_nine_anchors_respect_independent_margins():
    base = Rect(0, 0, 2, 1)
    expected = {
        "top_left": (1, 0.5),
        "top_center": (5.5, 0.5),
        "top_right": (10, 0.5),
        "middle_left": (1, 3.25),
        "middle_center": (5.5, 3.25),
        "middle_right": (10, 3.25),
        "bottom_left": (1, 6),
        "bottom_center": (5.5, 6),
        "bottom_right": (10, 6),
    }
    for anchor, (left, top) in expected.items():
        rect = resolve_block_rect(
            base,
            block_id="body",
            default_anchor="middle_center",
            style=_style(
                "body",
                anchor,
                top=0.5,
                right=1.333,
                bottom=0.5,
                left=1,
            ),
            slide_width=13.333,
            slide_height=7.5,
        )
        assert rect.left == left
        assert rect.top == top


def test_oversized_margins_keep_a_positive_box_inside_slide():
    rect = resolve_block_rect(
        Rect(0, 0, 8, 4),
        block_id="body",
        default_anchor="bottom_right",
        style=_style("body", "bottom_right", top=20, right=20, bottom=20, left=20),
        slide_width=10,
        slide_height=7.5,
    )
    assert rect.width > 0
    assert rect.height > 0
    assert rect.left >= 0
    assert rect.top >= 0
    assert rect.left + rect.width <= 10
    assert rect.top + rect.height <= 7.5


def test_explicit_zero_block_margins_reach_slide_edges():
    rect = resolve_block_rect(
        Rect(0.8, 2.4, 4, 2),
        block_id="title",
        default_anchor="middle_center",
        style=_style("title", "top_left", top=0, right=0, bottom=0, left=0),
        slide_width=13.333,
        slide_height=7.5,
    )
    assert rect.left == 0
    assert rect.top == 0


def test_layout_model_supports_extensible_block_ids():
    section = CoverSection(
        style=SectionStyle(
            blocks={
                "future_page_number": TextBlockStyle(
                    layout=BlockLayout(
                        anchor="bottom_right",
                        margin=EdgeInsets(right=0.4, bottom=0.3),
                    )
                )
            }
        )
    )
    restored = CoverSection.model_validate_json(section.model_dump_json())
    block = restored.style.blocks["future_page_number"].layout
    assert block.anchor == "bottom_right"
    assert block.margin.right == 0.4


def test_pptx_uses_relative_block_position_on_standard_slide(tmp_path):
    out = tmp_path / "layout.pptx"
    build_pptx(
        [
            SlideModel(
                kind="cover",
                section_id="s1",
                section_type="cover",
                title="标题",
                style=_style(
                    "title",
                    "bottom_right",
                    top=0.5,
                    right=0.5,
                    bottom=0.5,
                    left=0.5,
                ),
            )
        ],
        out,
        slide_size=SlideSize.STANDARD,
    )
    shape = Presentation(out).slides[0].shapes[0]
    assert round(shape.left.inches, 3) == 1.1
    assert round(shape.top.inches, 3) == 5.0
    assert round(shape.width.inches, 3) == 8.4
    assert round(shape.height.inches, 3) == 2.0


def test_pptx_uses_block_font_and_vertical_alignment(tmp_path):
    out = tmp_path / "block-font.pptx"
    build_pptx(
        [
            SlideModel(
                kind="cover",
                section_id="s1",
                section_type="cover",
                title="标题",
                style={
                    "title": {"font_size": 54, "color": "#111111"},
                    "blocks": {
                        "title": {
                            "text": {
                                "font_size": 31,
                                "color": "#123456",
                                "vertical_align": "bottom",
                            }
                        }
                    },
                },
            )
        ],
        out,
    )
    shape = Presentation(out).slides[0].shapes[0]
    run = shape.text_frame.paragraphs[0].runs[0]
    assert run.font.size.pt == 31
    assert str(run.font.color.rgb) == "123456"
    assert shape.text_frame.vertical_anchor == 4


def test_block_font_override_does_not_leak_to_sibling_body_role(tmp_path):
    out = tmp_path / "independent-block-fonts.pptx"
    build_pptx(
        [
            SlideModel(
                kind="cover",
                section_id="s1",
                section_type="cover",
                title="标题",
                subtitle="副标题",
                body="附加信息",
                style={
                    "body": {"font_size": 28, "color": "#222222"},
                    "blocks": {
                        "subtitle": {
                            "text": {"font_size": 19, "color": "#ABCDEF"}
                        }
                    },
                },
            )
        ],
        out,
    )
    shapes = Presentation(out).slides[0].shapes
    subtitle_run = shapes[1].text_frame.paragraphs[0].runs[0]
    extra_run = shapes[2].text_frame.paragraphs[0].runs[0]
    assert subtitle_run.font.size.pt == 19
    assert str(subtitle_run.font.color.rgb) == "ABCDEF"
    assert extra_run.font.size.pt == 28
    assert str(extra_run.font.color.rgb) == "222222"
