"""Tests for editable font styling in exported PPTX files."""
from __future__ import annotations

import zipfile

from pptx import Presentation

from app.pptx.builder import build_pptx
from app.services.generation import SlideModel


def test_pptx_exports_italic_underline_color_and_highlight(tmp_path):
    out = tmp_path / "font-style.pptx"
    build_pptx(
        [
            SlideModel(
                kind="media",
                section_id="s1",
                section_type="media",
                body="字体样式",
                style={
                    "body": {
                        "font_family": "Microsoft YaHei",
                        "font_size": 36,
                        "color": "#123456",
                        "bold": False,
                        "italic": True,
                        "underline": True,
                        "highlight_color": "#FFF200",
                        "align": "center",
                    }
                },
            )
        ],
        out,
    )

    presentation = Presentation(out)
    run = presentation.slides[0].shapes[0].text_frame.paragraphs[0].runs[0]
    assert run.font.bold is False
    assert run.font.italic is True
    assert run.font.underline is True
    assert str(run.font.color.rgb) == "123456"

    with zipfile.ZipFile(out) as archive:
        slide_xml = archive.read("ppt/slides/slide1.xml").decode("utf-8")
    assert '<a:highlight><a:srgbClr val="FFF200"/></a:highlight>' in slide_xml


def test_pptx_exports_scripture_verse_numbers_as_superscript_runs(tmp_path):
    out = tmp_path / "scripture-superscript.pptx"
    build_pptx(
        [
            SlideModel(
                kind="scripture",
                section_id="s1",
                section_type="scripture",
                body="1 起初神创造天地 2 地是空虚混沌",
                rich_body=[
                    [
                        {"text": "1", "superscript": True},
                        {"text": " 起初神创造天地"},
                        {"text": " "},
                        {"text": "2", "superscript": True},
                        {"text": " 地是空虚混沌"},
                    ]
                ],
                style={"body": {"font_size": 32}},
            )
        ],
        out,
    )

    with zipfile.ZipFile(out) as archive:
        slide_xml = archive.read("ppt/slides/slide1.xml").decode("utf-8")
    assert slide_xml.count('baseline="55000"') == 2
